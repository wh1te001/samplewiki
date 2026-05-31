from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Цвета проекта ──
NAVY = RGBColor(0x29, 0x28, 0x27)       # #292827
ACCENT = RGBColor(0xCC, 0x00, 0x00)     # #cc0000
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x50, 0x50, 0x50)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x88)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
SOFT_BG = RGBColor(0xF8, 0xF8, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, text="", font_size=12, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = align
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items, font_size=16, color=DARK_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
        # Bullet via prefix
        p.text = f"\u2022  {item}"
    return txBox


def add_placeholder_box(slide, left, top, width, height, label="[скриншот / схема]"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    shape.line.color.rgb = MEDIUM_GRAY
    shape.line.width = Pt(1.5)
    shape.line.dash_style = 2  # dash
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = MEDIUM_GRAY
    p.font.italic = True
    # Vertically center
    tf.paragraphs[0].space_before = Pt(0)
    shape.text_frame.auto_size = None
    return shape


def accent_bar(slide, top, width=Inches(0.08)):
    add_shape(slide, Inches(0), top, width, Inches(0.06), ACCENT)


def slide_header(slide, title, subtitle=None):
    # Top accent line
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), ACCENT)
    # Title
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), title,
                font_size=30, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4), subtitle,
                    font_size=15, color=MEDIUM_GRAY, bold=False)
    # Separator line
    add_shape(slide, Inches(0.8), Inches(1.35), Inches(2), Inches(0.04), ACCENT)


# ═══════════════════════════════════════════════
# СЛАЙД 1 — ТИТУЛЬНЫЙ
# ═══════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, NAVY)

# Accent stripe at top
add_shape(slide1, Inches(0), Inches(0), W, Inches(0.08), ACCENT)

# Big title
add_textbox(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "SampleWiki", font_size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
            "Каталог музыкальных сэмплов", font_size=28, color=RGBColor(0xCC, 0x66, 0x66), bold=False, align=PP_ALIGN.CENTER)

# Separator
add_shape(slide1, Inches(5.5), Inches(3.6), Inches(2.333), Inches(0.04), ACCENT)

add_textbox(slide1, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
            "Предзащита дипломного проекта", font_size=20, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slide1, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
            "ФИО, Группа, Руководитель", font_size=18, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

# Bottom line
add_shape(slide1, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT)


# ═══════════════════════════════════════════════
# СЛАЙД 2 — ПРОБЛЕМАТИКА
# ═══════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)
slide_header(slide2, "Проблематика и актуальность", "Почему SampleWiki?")

# Left column — Whosampled
add_textbox(slide2, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.4),
            "Whosampled.com (оригинал)", font_size=20, color=NAVY, bold=True)
problems = [
    "Долгая модерация — правки проверяются неделями",
    "Закрытый API, нет интеграций",
    "Нет поддержки Rutube и VK Video",
    "Платный доступ к расширенным функциям",
]
add_bullet_textbox(slide2, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5), problems,
                   font_size=15, color=DARK_GRAY)

# Placeholder for whosampled screenshot
add_placeholder_box(slide2, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.2),
                    "[Скриншот whosampled.com]")

# Right column — SampleWiki
add_textbox(slide2, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
            "SampleWiki (наш проект)", font_size=20, color=ACCENT, bold=True)
solutions = [
    "Российский аналог whosampled",
    "Формат вики — каждый может редактировать",
    "Встроенная поддержка Rutube и VK Video",
    "Открытый REST API для интеграций",
]
add_bullet_textbox(slide2, Inches(7), Inches(2.2), Inches(5.5), Inches(2.5), solutions,
                   font_size=15, color=DARK_GRAY)

# Placeholder for SampleWiki screenshot
add_placeholder_box(slide2, Inches(7), Inches(4.8), Inches(5.5), Inches(2.2),
                    "[Скриншот SampleWiki — главная страница]")


# ═══════════════════════════════════════════════
# СЛАЙД 3 — ЗАДАЧИ
# ═══════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)
slide_header(slide3, "Задачи проекта", "Что было реализовано")

tasks = [
    "Спроектировать модель данных (Artist → Album → Track → Sample)",
    "Разработать REST API на ASP.NET Core 8",
    "Реализовать аутентификацию через JWT в httpOnly cookie",
    "Создать MPA-фронтенд в стиле whosampled.com",
    "Встроить видеоплееры VK Video и Rutube",
    "Реализовать поиск с автокомплитом",
    "Добавить reverse samples (отображение \"использован в\")",
    "Обеспечить аудит изменений через систему Revision",
]
add_bullet_textbox(slide3, Inches(0.8), Inches(1.7), Inches(7.5), Inches(5), tasks,
                   font_size=16, color=DARK_GRAY)

# Placeholder for diagram
add_placeholder_box(slide3, Inches(8.8), Inches(1.7), Inches(4), Inches(5),
                    "[Диаграмма —\nсхема проекта]")


# ═══════════════════════════════════════════════
# СЛАЙД 4 — СТЕК ТЕХНОЛОГИЙ
# ═══════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)
slide_header(slide4, "Технологический стек")

techs = [
    ("Backend API", "ASP.NET Core (.NET 8)", "REST API, Swagger, EF Core"),
    ("База данных", "MySQL 8 + Entity Framework Core", "Pomelo.EntityFrameworkCore.MySql"),
    ("Аутентификация", "JWT (HMAC-SHA256)", "httpOnly cookie, BCrypt"),
    ("Frontend", "HTML + Vanilla JS + CSS", "MPA, Fetch API, modules"),
    ("Статический сервер", "Python http.server", "порт 8000"),
    ("Rate Limiting", "ASP.NET FixedWindow", "10 запросов/мин на IP"),
]

# Table header
y_start = Inches(1.7)
row_h = Inches(0.55)
col_w = [Inches(2.5), Inches(4), Inches(5)]

headers = ["Компонент", "Технология", "Детали"]
for ci, h in enumerate(headers):
    x = Inches(0.8) + sum(col_w[:ci]) + Inches(ci * 0.15)
    w = col_w[ci]
    add_shape(slide4, Inches(0.8) + sum(col_w[:ci]) + Inches(ci * 0.15),
              y_start, w, row_h, NAVY, h, font_size=15, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for ri, (comp, tech, detail) in enumerate(techs):
    ry = y_start + row_h + Inches(ri * (row_h + Inches(0.08)))
    bg = RGBColor(0xF5, 0xF5, 0xF5) if ri % 2 == 0 else WHITE
    vals = [comp, tech, detail]
    for ci, v in enumerate(vals):
        add_shape(slide4, Inches(0.8) + sum(col_w[:ci]) + Inches(ci * 0.15),
                  ry, col_w[ci], row_h, bg, v, font_size=14, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# СЛАЙД 5 — АРХИТЕКТУРА
# ═══════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)
slide_header(slide5, "Архитектура приложения", "MPA + CORS — два независимых сервера")

# Architecture diagram — boxes with arrows
# Browser box
add_shape(slide5, Inches(0.8), Inches(1.8), Inches(2.5), Inches(1.5), NAVY,
          "Браузер\n(Chrome / Firefox)", font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Arrow right
add_textbox(slide5, Inches(3.3), Inches(2.3), Inches(1.5), Inches(0.5),
            "GET  →", font_size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Python server
add_shape(slide5, Inches(4.5), Inches(1.8), Inches(2.8), Inches(1.5), RGBColor(0x50, 0x50, 0x50),
          "Python http.server\n:8000\nСтатика (HTML/JS/CSS)", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Arrow right — fetch
add_textbox(slide5, Inches(7.3), Inches(2.3), Inches(1.5), Inches(0.5),
            "fetch  →", font_size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ASP.NET box
add_shape(slide5, Inches(8.5), Inches(1.8), Inches(2.8), Inches(1.5), ACCENT,
          "ASP.NET Core\n:5000\nREST API", font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Arrow right
add_textbox(slide5, Inches(11.3), Inches(2.3), Inches(1), Inches(0.5),
            " → SQL", font_size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# MySQL
add_shape(slide5, Inches(12), Inches(1.8), Inches(1), Inches(1.5), NAVY,
          "MySQL", font_size=13, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Key points below
add_textbox(slide5, Inches(0.8), Inches(3.8), Inches(12), Inches(0.4),
            "Ключевые особенности архитектуры:", font_size=18, color=NAVY, bold=True)
arch_points = [
    "Два сервера: статический (Python :8000) и API (ASP.NET :5000)",
    "Кросс-доменные запросы через fetch с credentials: 'include'",
    "CORS настроен только для localhost:8000 с AllowCredentials",
    "Фронтенд — MPA (Multi Page Application), 9 HTML-страниц",
    "Все ссылки через ?id= в query string",
]
add_bullet_textbox(slide5, Inches(0.8), Inches(4.3), Inches(11), Inches(3), arch_points,
                   font_size=15, color=DARK_GRAY)

# Placeholder for architecture scheme
add_placeholder_box(slide5, Inches(7), Inches(4.5), Inches(5.5), Inches(2.5),
                    "[Полная схема архитектуры]")


# ═══════════════════════════════════════════════
# СЛАЙД 6 — МОДЕЛЬ ДАННЫХ
# ═══════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, WHITE)
slide_header(slide6, "Модель данных", "7 сущностей, двунаправленные связи сэмплов")

# ER diagram as boxes with relationships
entities = [
    ("Artist", Inches(0.8), Inches(2.0), Inches(2.5), Inches(0.9)),
    ("Album", Inches(4.5), Inches(2.0), Inches(2.5), Inches(0.9)),
    ("Track", Inches(8.5), Inches(2.0), Inches(2.5), Inches(0.9)),
    ("Sample", Inches(10.5), Inches(4.0), Inches(2.3), Inches(0.9)),
]

other_entities = [
    ("User", Inches(0.8), Inches(5.5), Inches(2.5), Inches(0.8)),
    ("Artwork", Inches(4.5), Inches(5.5), Inches(2.5), Inches(0.8)),
    ("Revision", Inches(8.5), Inches(5.5), Inches(2.5), Inches(0.8)),
]

for name, x, y, w, h in entities:
    add_shape(slide6, x, y, w, h, ACCENT if name == "Sample" else NAVY,
              name, font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for name, x, y, w, h in other_entities:
    add_shape(slide6, x, y, w, h, DARK_GRAY,
              name, font_size=16, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Relationship arrows (text-based)
rels = [
    (Inches(3.3), Inches(2.3), "1—N", Inches(4.5)),
    (Inches(7), Inches(2.3), "1—N", Inches(8.5)),
]
for x, y, label, x2 in rels:
    add_textbox(slide6, x, y, Inches(0.6), Inches(0.4), label, font_size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    # Simple arrow line
    add_shape(slide6, x + Inches(0.6), y + Inches(0.15), x2 - x - Inches(0.6), Inches(0.02), ACCENT)

# SampledBy arrow (reverse)
add_textbox(slide6, Inches(8.8), Inches(3.0), Inches(1.5), Inches(0.4),
            "← SampledBy", font_size=12, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_shape(slide6, Inches(9), Inches(3.4), Inches(1.8), Inches(0.02), ACCENT)

# Fields list
add_textbox(slide6, Inches(0.8), Inches(6.5), Inches(12), Inches(0.6),
            "Artist (Name, Description, WikiLink)  →  Album (Title, ReleaseYear)  →  Track (Title, Genre, #)  →  Sample (Type, StartTimeSeconds)",
            font_size=12, color=DARK_GRAY, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════
# СЛАЙД 7 — REST API
# ═══════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, WHITE)
slide_header(slide7, "REST API", "9 контроллеров, 30+ эндпоинтов")

api_groups = [
    ("Auth", "POST /api/auth/register|login|logout\nGET /api/auth/me", "httpOnly cookie, rate limited"),
    ("Artists", "GET /api/artists[/{id}]\nPOST /api/artists", "с пагинацией, TrackCount"),
    ("Tracks", "GET /api/tracks[/{id}]\nPOST/PUT /api/tracks/{id}", "с Samples, SampledBy, Album, Artist"),
    ("Albums", "GET /api/albums/{id}", "с треками, обложками, исполнителем"),
    ("Samples", "GET /api/samples[/{id}|/track/{id}]\nPOST/PUT /api/samples", "с деталями треков"),
    ("Search", "GET /api/search/artists?q=\nGET /api/search/tracks?q=&artistId=", "автокомплит"),
    ("Submit", "POST /api/submit", "создание сэмпла в 1 транзакции"),
    ("Upload", "POST /api/upload", "загрузка изображений"),
    ("Revisions", "GET /api/revisions[/track/{id}|/user/{id}]", "аудит изменений"),
]

# Table
y_s = Inches(1.7)
rh = Inches(0.55)
cw = [Inches(1.5), Inches(6), Inches(4.3)]

headers7 = ["Группа", "Эндпоинты", "Особенности"]
for ci, h in enumerate(headers7):
    add_shape(slide7, Inches(0.8) + sum(cw[:ci]) + Inches(ci * 0.1),
              y_s, cw[ci], rh, NAVY, h, font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for ri, (grp, eps, note) in enumerate(api_groups):
    ry = y_s + rh + Inches(ri * (rh + Inches(0.05)))
    bg = RGBColor(0xF5, 0xF5, 0xF5) if ri % 2 == 0 else WHITE
    vals = [grp, eps, note]
    for ci, v in enumerate(vals):
        fs = 13 if ci > 0 else 14
        add_shape(slide7, Inches(0.8) + sum(cw[:ci]) + Inches(ci * 0.1),
                  ry, cw[ci], rh, bg, v, font_size=fs, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Placeholder
add_placeholder_box(slide7, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5),
                    "[Скриншот Swagger UI]")


# ═══════════════════════════════════════════════
# СЛАЙД 8 — FRONTEND
# ═══════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, WHITE)
slide_header(slide8, "Frontend", "MPA в стиле whosampled.com")

# Left: pages list
add_textbox(slide8, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4),
            "Страницы:", font_size=20, color=NAVY, bold=True)

pages = [
    "index.html — главная (список исполнителей)",
    "artist.html — детально об исполнителе",
    "album.html — альбом с треками и обложками",
    "track.html — трек с сэмплами и SampledBy",
    "sample.html — детально о сэмпле",
    "search.html — глобальный поиск",
    "auth.html / register.html — вход / регистрация",
    "submit.html — создание сэмпла (авторизованные)",
]
add_bullet_textbox(slide8, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), pages,
                   font_size=14, color=DARK_GRAY)

# Right: JS modules
add_textbox(slide8, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
            "JS модули:", font_size=20, color=NAVY, bold=True)

modules = [
    "api.js — fetch-обёртка с credentials: 'include'",
    "auth.js — register / login / logout / checkAuth",
    "render.js — отрисовка всех страниц",
    "utils.js — toast, форматирование, navbarSearch",
    "submit.js — автокомплит, загрузка, отправка",
]
add_bullet_textbox(slide8, Inches(7), Inches(2.2), Inches(5.5), Inches(3), modules,
                   font_size=14, color=DARK_GRAY)

# Design highlights
add_textbox(slide8, Inches(7), Inches(4.5), Inches(5.5), Inches(0.4),
            "Особенности дизайна:", font_size=18, color=NAVY, bold=True)
design = [
    "Цвета: #292827 (navbar), #cc0000 (акцент)",
    "Navbar: grid-layout (лого / поиск / auth)",
    "Breadcrumb навигация, Song Connections",
    "Inline-поиск на всех страницах",
]
add_bullet_textbox(slide8, Inches(7), Inches(5.0), Inches(5.5), Inches(2), design,
                   font_size=14, color=DARK_GRAY)

add_placeholder_box(slide8, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6),
                    "[Скриншоты страниц: track.html, album.html, index.html]")


# ═══════════════════════════════════════════════
# СЛАЙД 9 — КЛЮЧЕВЫЕ ВОЗМОЖНОСТИ
# ═══════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, WHITE)
slide_header(slide9, "Ключевые возможности")

features = [
    ("Embed-виджеты", "VK Video, Rutube, YouTube — встроенные плееры\nс jump-to-time по StartTimeSeconds", ACCENT),
    ("Reverse Samples", "Секция «Использован в» на странице трека —\nкакие треки засемплировали этот", RGBColor(0x2E, 0x7D, 0x32)),
    ("Автокомплит", "Поиск артистов и треков при создании сэмпла\nс debounce 300 мс", NAVY),
    ("Аудит Revision", "Полная история изменений каждого трека\nс хранением OldValues / NewValues в JSON", DARK_GRAY),
    ("Submit", "Создание артиста + альбома + трека + сэмпла\nв одной транзакции", RGBColor(0x15, 0x65, 0xC0)),
]

for i, (title, desc, color) in enumerate(features):
    x = Inches(0.8) + (i % 3) * Inches(4.1)
    y = Inches(1.8) + (i // 3) * Inches(2.7)
    # Card
    card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    card.line.width = Pt(1)
    # Top color bar
    add_shape(slide9, x + Inches(0.2), y + Inches(0.15), Inches(3.4), Inches(0.06), color)
    # Title
    add_textbox(slide9, x + Inches(0.3), y + Inches(0.4), Inches(3.2), Inches(0.5),
                title, font_size=17, color=color, bold=True)
    # Description
    add_textbox(slide9, x + Inches(0.3), y + Inches(0.9), Inches(3.2), Inches(1.2),
                desc, font_size=13, color=DARK_GRAY)

add_placeholder_box(slide9, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5),
                    "[Скриншоты: embed плеера, reverse samples, автокомплита, Swagger]")


# ═══════════════════════════════════════════════
# СЛАЙД 10 — БЕЗОПАСНОСТЬ
# ═══════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, WHITE)
slide_header(slide10, "Безопасность")

sec_items = [
    ("JWT в httpOnly cookie", "Токен недоступен из JavaScript — защита от XSS-атак"),
    ("SameSite=Strict", "Cookie не отправляется на сторонние сайты — защита от CSRF"),
    ("BCrypt хеширование", "Пароли никогда не хранятся в открытом виде"),
    ("Rate Limiting", "10 запросов в минуту на IP на AuthController"),
    ("Валидация данных", "Data Annotations на бэкенде + inline-валидация на фронтенде"),
    ("CORS", "Только origin http://localhost:8000 с AllowCredentials"),
    ("Блокировка аккаунтов", "Поле IsActive — деактивированные пользователи не впускаются"),
    ("Логирование", "ILogger: регистрация, вход, ошибки (Warn/Error)"),
]

y_pos = Inches(1.7)
for i, (title, desc) in enumerate(sec_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6.2)
    y = y_pos + row * Inches(1.35)
    # Icon circle
    circle = slide10.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.05), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    add_textbox(slide10, x + Inches(0.05), y + Inches(0.05), Inches(0.4), Inches(0.4),
                str(i + 1), font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide10, x + Inches(0.55), y, Inches(5.3), Inches(0.35),
                title, font_size=16, color=NAVY, bold=True)
    # Description
    add_textbox(slide10, x + Inches(0.55), y + Inches(0.35), Inches(5.3), Inches(0.6),
                desc, font_size=13, color=DARK_GRAY)


# ═══════════════════════════════════════════════
# СЛАЙД 11 — РЕЗУЛЬТАТЫ
# ═══════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide11, WHITE)
slide_header(slide11, "Результаты и дальнейшие планы")

# Left column — done
add_shape(slide11, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5), GREEN,
          "✓  Реализовано", font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

done_items = [
    "REST API: чтение/создание/обновление артистов, треков, сэмплов, альбомов",
    "Аутентификация (регистрация, вход, httpOnly cookie)",
    "Поиск с автокомплитом (/api/search)",
    "Reverse samples (двунаправленные связи сэмплов)",
    "Встраивание плееров (VK, Rutube, YouTube, SoundCloud)",
    "Создание сэмпла в одной транзакции (Submit)",
    "Загрузка изображений (Upload)",
    "История изменений (Revision / Audit Trail)",
    "MPA-фронтенд в стиле whosampled.com",
]
add_bullet_textbox(slide11, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.5), done_items,
                   font_size=14, color=DARK_GRAY)

# Right column — todo
add_shape(slide11, Inches(7), Inches(1.7), Inches(5.5), Inches(0.5), MEDIUM_GRAY,
          "○  В планах", font_size=18, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

todo_items = [
    "Удаление сущностей через UI (сейчас только через API)",
    "Админ-панель для модерации",
    "Пагинация больших списков",
    "Кэширование на фронтенде (Service Worker)",
    "Unit-тесты и E2E-тесты",
    "HTTPS для продакшена (Secure cookie)",
    "Полнотекстовый поиск (FULLTEXT index в MySQL)",
]
add_bullet_textbox(slide11, Inches(7), Inches(2.4), Inches(5.5), Inches(4.5), todo_items,
                   font_size=14, color=DARK_GRAY)


# ═══════════════════════════════════════════════
# СЛАЙД 12 — СПАСИБО
# ═══════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide12, NAVY)

add_shape(slide12, Inches(0), Inches(0), W, Inches(0.08), ACCENT)

add_textbox(slide12, Inches(1), Inches(2.5), Inches(11), Inches(1),
            "Спасибо за внимание!", font_size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_shape(slide12, Inches(5.5), Inches(3.6), Inches(2.333), Inches(0.04), ACCENT)

add_textbox(slide12, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
            "Вопросы?", font_size=28, color=RGBColor(0xCC, 0x66, 0x66), align=PP_ALIGN.CENTER)

add_textbox(slide12, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
            "GitHub: ссылка на репозиторий", font_size=16, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

add_shape(slide12, Inches(0), H - Inches(0.08), W, Inches(0.08), ACCENT)


# ── Сохраняем ──
output_path = "C:\\xampp\\htdocs\\sw3\\samplewiki\\SampleWiki_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
